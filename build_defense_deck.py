# -*- coding: utf-8 -*-
"""LogicSplat M.Sc. thesis defense deck — built from manuscript.tex (authoritative).
Every number traces to the manuscript; no estimates."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from lxml.etree import SubElement
import copy

# ---------------- design tokens ----------------
NAVY      = RGBColor(0x10, 0x2A, 0x43)   # primary dark
NAVY_2    = RGBColor(0x24, 0x3B, 0x53)   # panel on navy
NAVY_3    = RGBColor(0x33, 0x4E, 0x68)   # motif nodes
INK       = RGBColor(0x1F, 0x29, 0x33)   # headlines
BODY      = RGBColor(0x3E, 0x4C, 0x59)   # body text
MUTED     = RGBColor(0x7B, 0x87, 0x94)   # secondary
FAINT     = RGBColor(0x9F, 0xB3, 0xC8)   # muted bars / faint text
ACCENT    = RGBColor(0xC2, 0x41, 0x0C)   # burnt orange
ACCENT_D  = RGBColor(0xE8, 0x7A, 0x3E)   # accent on dark bg
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
CLOUD     = RGBColor(0xBC, 0xCC, 0xDC)   # light text on navy
PANEL     = RGBColor(0xF0, 0xF4, 0xF8)   # card fill
PANEL_LN  = RGBColor(0xD9, 0xE2, 0xEC)   # hairlines
BAR_MUTE  = RGBColor(0x9F, 0xB3, 0xC8)   # baseline bars
BAR_NAVY  = RGBColor(0x48, 0x65, 0x81)   # secondary "ours" bar

F    = "Segoe UI"
FSB  = "Segoe UI Semibold"
MONO = "Consolas"

SW, SH = 13.333, 7.5
ML = 0.62                    # side margin
CW = SW - 2 * ML             # content width = 12.093

FIG = r"C:\Users\Debjyoti\Desktop\LogicSplat\figures"
OUT = r"C:\Users\Debjyoti\Desktop\LogicSplat\LogicSplat_Defense.pptx"

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

TOTAL = 22
_n = [0]

# ---------------- helpers ----------------
def slide_new(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    _n[0] += 1
    return s

def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    return tf

def para(tf, runs, size=12, color=BODY, bold=False, align=PP_ALIGN.LEFT,
         before=0, after=0, leading=1.06, font=F, first=False, italic=False):
    """runs: str or list of (text, overrides-dict)"""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    p.line_spacing = leading
    if isinstance(runs, str):
        runs = [(runs, {})]
    for text, ov in runs:
        r = p.add_run()
        r.text = text
        r.font.name = ov.get("font", font)
        r.font.size = Pt(ov.get("size", size))
        r.font.bold = ov.get("bold", bold)
        r.font.italic = ov.get("italic", italic)
        r.font.color.rgb = ov.get("color", color)
    return p

def rect(s, x, y, w, h, fill, line=None, line_w=0.75, radius=None):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    shp = s.shapes.add_shape(shp_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius is not None:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def hline(s, x, y, w, color=PANEL_LN, weight=0.75):
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y),
                                Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    return ln

def vline(s, x, y, h, color=PANEL_LN, weight=0.75):
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y),
                                Inches(x), Inches(y + h))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    return ln

def footer(s):
    n = _n[0]
    hline(s, ML, 7.08, CW, PANEL_LN, 0.75)
    tf = tb(s, ML, 7.14, 8.0, 0.26)
    para(tf, "LogicSplat  ·  M.Sc. Thesis Defense  ·  Symbiosis Institute of Geoinformatics",
         size=8.5, color=MUTED, first=True)
    tf = tb(s, SW - ML - 1.5, 7.14, 1.5, 0.26)
    para(tf, f"{n:02d} / {TOTAL}", size=8.5, color=MUTED, align=PP_ALIGN.RIGHT, first=True)

def header(s, kicker, headline, headline_size=25):
    rect(s, ML, 0.44, 0.42, 0.055, ACCENT)
    tf = tb(s, ML, 0.58, CW, 0.3)
    para(tf, kicker.upper(), size=11, color=ACCENT, bold=True, first=True, font=FSB)
    tf = tb(s, ML, 0.90, CW, 1.0)
    para(tf, headline, size=headline_size, color=INK, bold=True, first=True,
         leading=1.04, font=FSB)

def std_slide(kicker, headline, headline_size=25):
    s = slide_new()
    header(s, kicker, headline, headline_size)
    footer(s)
    return s

def bar_row(s, x, y, label, sub, value, vmax, bar_area_w, fill,
            label_w=3.0, bar_h=0.30, value_str=None):
    """One horizontal bar: right-aligned label | bar | value in ink."""
    tf = tb(s, x, y - 0.02, label_w, 0.55)
    para(tf, label, size=11.5, color=INK, bold=True, align=PP_ALIGN.RIGHT, first=True)
    if sub:
        para(tf, sub, size=8, color=MUTED, align=PP_ALIGN.RIGHT, before=1)
    bx = x + label_w + 0.18
    bw = max(0.06, (value / vmax) * bar_area_w)
    rect(s, bx, y, bw, bar_h, fill, radius=0.5)
    tf = tb(s, bx + bw + 0.12, y - 0.035, 1.2, 0.36)
    para(tf, value_str or f"{value:.1f}%", size=12.5, color=INK, bold=True,
         first=True, font=FSB)
    return bx

def pic(s, path, x, y, w=None, h=None):
    kw = {}
    if w is not None:
        kw["width"] = Inches(w)
    if h is not None:
        kw["height"] = Inches(h)
    return s.shapes.add_picture(path, Inches(x), Inches(y), **kw)

def motif(s, pts, edges, node_fill=NAVY_3, edge_color=NAVY_2, accent_idx=None, r=0.11):
    for a, b in edges:
        x1, y1 = pts[a]; x2, y2 = pts[b]
        ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                                    Inches(x2), Inches(y2))
        ln.line.color.rgb = edge_color
        ln.line.width = Pt(1.4)
        ln.shadow.inherit = False
    for i, (px, py) in enumerate(pts):
        fill = ACCENT_D if i == accent_idx else node_fill
        sz = r * (1.55 if i == accent_idx else 1.0)
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px - sz), Inches(py - sz),
                               Inches(2 * sz), Inches(2 * sz))
        c.fill.solid(); c.fill.fore_color.rgb = fill
        c.line.fill.background(); c.shadow.inherit = False

# ======================================================================
# SLIDE 1 — TITLE
# ======================================================================
s = slide_new(bg=NAVY)

# -- title graphic: a soft cloud of translucent Gaussian ellipses (a "splat") --
import random as _rnd
_rnd.seed(11)

def gauss_blob(x, y, w, h, color, alpha_pct, rot):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - w / 2), Inches(y - h / 2),
                             Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    sF = shp.fill._xPr.find(qn('a:solidFill'))
    clr = sF.find(qn('a:srgbClr'))
    SubElement(clr, qn('a:alpha'), {'val': str(int(alpha_pct * 1000))})
    shp.line.fill.background()
    shp.shadow.inherit = False
    shp.rotation = rot
    return shp

B2 = RGBColor(0x2E, 0x4C, 0x6B)
B3 = RGBColor(0x3A, 0x5F, 0x85)
B4 = RGBColor(0x4A, 0x74, 0x9E)
B5 = RGBColor(0x6B, 0x94, 0xBE)

# layered translucent gaussians, denser toward the cloud centre
palette = [B2, B2, B3, B3, B3, B4, B4, B5]
for _ in range(34):
    bx = _rnd.gauss(11.35, 1.05)
    by = min(6.7, max(0.7, _rnd.gauss(3.55, 1.45)))
    if bx < 9.55:
        bx = 9.55 + (9.55 - bx) * 0.35
    bw = min(1.75, max(0.28, _rnd.lognormvariate(-0.45, 0.55)))
    bh = bw / _rnd.uniform(1.5, 2.9)
    col = _rnd.choice(palette)
    al = _rnd.uniform(11, 24) if bw > 0.85 else _rnd.uniform(20, 38)
    gauss_blob(bx, by, bw, bh, col, al, _rnd.uniform(0, 180))

# small bright amber gaussians — glowing accents, not washes
gauss_blob(11.95, 2.60, 0.62, 0.28, ACCENT_D, 58, 24)
gauss_blob(10.55, 4.75, 0.44, 0.20, ACCENT_D, 48, 140)

# fine particles, like individual splat points
for _ in range(26):
    px_ = _rnd.gauss(11.35, 1.25)
    py_ = _rnd.gauss(3.55, 1.6)
    if px_ < 9.35 or px_ > 13.1 or py_ < 0.6 or py_ > 6.9:
        continue
    d = _rnd.uniform(0.035, 0.085)
    col = _rnd.choice([B4, B5, B5, CLOUD])
    gauss_blob(px_, py_, d, d, col, _rnd.uniform(40, 75), 0)
# three amber sparks
for sx, sy, sd in [(12.55, 1.85, 0.07), (10.15, 2.95, 0.055), (11.7, 5.6, 0.065)]:
    gauss_blob(sx, sy, sd, sd, ACCENT_D, 85, 0)

tf = tb(s, ML, 1.30, 9.0, 0.35)
para(tf, "M.SC. THESIS DEFENSE  ·  JULY 2026", size=12, color=ACCENT_D,
     bold=True, first=True, font=FSB)
rect(s, ML, 1.78, 0.55, 0.06, ACCENT_D)

tf = tb(s, ML, 2.05, 9.3, 1.2)
para(tf, "LogicSplat", size=56, color=WHITE, bold=True, first=True, font=FSB)
tf = tb(s, ML, 3.22, 8.7, 1.1)
para(tf, "Neuro-Symbolic 3D Scene Graph Generation via\nGeometric Kolmogorov–Arnold Networks",
     size=20, color=CLOUD, first=True, leading=1.15)

hline(s, ML, 5.05, 5.6, NAVY_3, 1.0)
tf = tb(s, ML, 5.30, 9.0, 1.6)
para(tf, [("Debjyoti Sengupta", {"bold": True, "color": WHITE, "size": 16, "font": FSB}),
          ("    PRN 24070243013", {"color": FAINT, "size": 12.5})], first=True, after=8)
para(tf, [("Internal Guide:  ", {"color": FAINT, "size": 12}),
          ("Mr. Sahil Shah", {"color": CLOUD, "size": 12, "bold": True}),
          ("  ·  Symbiosis Institute of Geoinformatics", {"color": FAINT, "size": 12})], after=4)
para(tf, "Symbiosis Institute of Geoinformatics, Symbiosis International (Deemed University), Pune",
     size=12, color=FAINT)

# institution logo, top-right: white plate so the red emblem sits on its own
# ground rather than clashing with the navy (drawn after the gaussian cloud,
# so it reads as the topmost, most deliberate layer)
LOGO_W, LOGO_H = 0.80, 0.80 * (197 / 235)   # true aspect of figures/sig_logo.png
PLATE_W, PLATE_H = 1.10, 0.94
plate_x = SW - ML - PLATE_W
plate_y = 0.50
rect(s, plate_x, plate_y, PLATE_W, PLATE_H, WHITE, radius=0.10)
pic(s, FIG + r"\sig_logo.png",
    plate_x + (PLATE_W - LOGO_W) / 2, plate_y + (PLATE_H - LOGO_H) / 2, w=LOGO_W)

# ======================================================================
# SLIDE 2 — CONTEXT / HOOK
# ======================================================================
s = std_slide("Context",
              "Machines that act in the world need relations, not just objects.")
tf = tb(s, ML, 2.05, 5.75, 4.6)
para(tf, [("A robot setting a table must know not merely that a cup exists, but that it ",
           {}),
          ("rests on", {"bold": True, "color": INK}),
          (" the table. An augmented-reality assistant must know that a router sits ", {}),
          ("behind", {"bold": True, "color": INK}),
          (" a laptop, not beside it.", {})],
     size=13, color=BODY, first=True, leading=1.28, after=12)
para(tf, [("A ", {}),
          ("scene graph", {"bold": True, "color": INK}),
          (" captures exactly this structure: object instances as nodes, directed spatial "
           "relations as edges.", {})],
     size=13, color=BODY, leading=1.28, after=12)
para(tf, "3D Gaussian Splatting has made the geometry itself inexpensive: a photorealistic, "
         "room-scale reconstruction from a short smartphone video, in five to fifteen minutes, "
         "on consumer hardware.",
     size=13, color=BODY, leading=1.28)
rect(s, ML, 5.55, 5.75, 1.06, PANEL, PANEL_LN, radius=0.08)
tf = tb(s, ML + 0.28, 5.55, 5.75 - 0.56, 1.06, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("The question this thesis asks:  ", {"bold": True, "color": ACCENT, "font": FSB}),
          ("is that geometry, on its own, sufficient to recover the relational structure of "
           "a scene?", {"color": INK})],
     size=12.5, first=True, leading=1.2)
hk_h = 1.65
hk_w1 = hk_h * 1.638
hk_w2 = hk_h * 1.4184
hk_x = 7.0
hk_y = 2.95
pic(s, FIG + r"\fig_input_photo.png", hk_x, hk_y, h=hk_h)
tf = tb(s, hk_x + hk_w1 + 0.06, hk_y, 0.5, hk_h, anchor=MSO_ANCHOR.MIDDLE)
para(tf, "→", size=20, color=ACCENT, bold=True, align=PP_ALIGN.CENTER, first=True)
pic(s, FIG + r"\fig_splat_reconstruction.png", hk_x + hk_w1 + 0.62, hk_y, h=hk_h)
tf = tb(s, hk_x, hk_y + hk_h + 0.16, hk_w1 + hk_w2 + 0.62, 0.6)
para(tf, "A handheld video becomes a photorealistic 3D reconstruction in minutes — "
         "reconstruction is no longer the obstacle; understanding it is.",
     size=10, color=MUTED, first=True, leading=1.2, align=PP_ALIGN.CENTER, italic=True)

# ======================================================================
# SLIDE 3 — MOTIVATION
# ======================================================================
s = std_slide("Motivation",
              "Three obstacles separate a Gaussian splat from a reliable scene graph.")
cards = [
    ("01", "Semantic grounding is expensive",
     [("Existing systems — GaussianGraph, ReLaGS, ConceptGraphs, Open3DSG, RelationField — treat knowing "
       "what an object is as a prerequisite for reasoning about how objects relate.", {}),
      ],
     "Hundreds of millions to billions of parameters; per-scene optimisation; in several cases over ten minutes of inference — for relations that are fully determined by centroids and bounding volumes."),
    ("02", "3DSSG is severely under-annotated",
     [("Human labellers mark only ", {}),
      ("80–120", {"bold": True, "color": INK}),
      (" directed pairs per scene out of roughly ", {}),
      ("870", {"bold": True, "color": INK}),
      (" geometrically valid pairs — ≈90% of pairs are unannotated.", {})],
     "Standard supervised training treats them as negatives, injecting a false-negative signal that suppresses recall for relations geometry could decide with certainty."),
    ("03", "Predicted graphs are logically inconsistent",
     [("Relations are predicted independently per pair: ", {}),
      ("higher_than(A,B)", {"font": MONO, "size": 10.5, "color": INK}),
      (" and ", {}),
      ("higher_than(B,A)", {"font": MONO, "size": 10.5, "color": INK}),
      (" can coexist; ", {}),
      ("on_top_of(A,B)", {"font": MONO, "size": 10.5, "color": INK}),
      (" can appear without ", {}),
      ("under(B,A)", {"font": MONO, "size": 10.5, "color": INK}),
      (".", {})],
     "No existing system applies a formal, parameter-free mechanism to resolve such contradictions — despite most constraints being universal physical facts."),
]
cx = ML
for num, title, body_runs, tail in cards:
    rect(s, cx, 1.90, 3.83, 4.85, PANEL, PANEL_LN, radius=0.05)
    tf = tb(s, cx + 0.28, 2.16, 3.27, 4.4)
    para(tf, num, size=24, color=FAINT, bold=True, first=True, font=FSB, after=4)
    para(tf, title, size=14.5, color=INK, bold=True, font=FSB, after=8, leading=1.05)
    para(tf, body_runs, size=11.5, color=BODY, after=8, leading=1.16)
    para(tf, tail, size=10.5, color=MUTED, leading=1.16)
    cx += 3.83 + 0.30

# ======================================================================
# SLIDE 3 — OBJECTIVES
# ======================================================================
s = std_slide("Objectives",
              "Four objectives, each answered by a controlled experiment.")
objs = [
    ("O1", "Geometry-only relation prediction",
     "Predict spatial relations directly from Gaussian-splat geometry using GeoKANRelationGNN, built on "
     "Geometric Kolmogorov–Arnold Networks — without CLIP, SAM, a large language model, or per-scene "
     "optimisation at inference."),
    ("O2", "Repair the training signal",
     "Address the severe annotation sparsity of 3DSSG through a deterministic geometric rule that derives "
     "pseudo-labels for unannotated object pairs, rather than treating them as false negatives."),
    ("O3", "Guarantee logical consistency",
     "Resolve logical contradictions and complete missing relations in the predicted scene graph through "
     "SceneGraphRepair — a zero-parameter symbolic component using deterministic fixed-point iteration."),
    ("O4", "Test the learnable metric fairly",
     "Test whether a learnable-metric architecture (GeoKAN) provides a measurable advantage over a "
     "matched-capacity multilayer perceptron — both in-distribution and under a change of domain."),
]
oy = 1.95
for tag, title, body in objs:
    chip = rect(s, ML, oy + 0.06, 0.62, 0.62, NAVY, radius=0.18)
    tf = chip.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    para(tf, tag, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER, first=True, font=FSB)
    tf = tb(s, ML + 0.95, oy, CW - 0.95, 1.1)
    para(tf, title, size=14.5, color=INK, bold=True, first=True, font=FSB, after=3)
    para(tf, body, size=11.5, color=BODY, leading=1.14)
    oy += 1.22
    if tag != "O4":
        hline(s, ML + 0.95, oy - 0.17, CW - 0.95, PANEL_LN, 0.75)

# ======================================================================
# SLIDE 4 — RELATED WORK
# ======================================================================
s = std_slide("Related Work",
              "Existing systems reach relations only through foundation models or per-scene optimisation.")
cols = [("System", ML, 2.05), ("Approach", 2.85, 4.75), ("Needs at inference", 7.78, 2.55),
        ("Recall@5", 10.55, 2.15)]
rows = [
    ("ReLaGS", "CVPR 2026", "CLIP + SAM Gaussian features; GPT-4o relation annotation; GNN-based reasoning",
     "CLIP + SAM (~1.5 B params)", "87.0%", "3DSSG"),
    ("RelationField", "CVPR 2025", "GPT-4o relation knowledge distilled into a NeRF optimised per scene",
     "~60–90 min per-scene optimisation (A100)", "82.0%", "3DSSG"),
    ("ConceptGraphs", "2023", "SAM + CLIP object features, LLaVA captions, relations inferred by GPT-4",
     "Large language model", "79.0%", "3DSSG"),
    ("Open3DSG", "CVPR 2024", "3D backbone co-embedded with CLIP, OpenSeg, InstructBLIP / Vicuna-7B",
     "CLIP + LLM", "65.0%", "3DSSG"),
    ("GaussianGraph", "2025", "CLIP, SAM2, LLaVA-1.6, Grounding DINO, plus a 3D relation-correction module",
     "Four foundation models", "63.2%", "LERF, positional"),
]
ty = 1.95
for name, x, w in cols:
    tf = tb(s, x, ty, w, 0.3)
    para(tf, name.upper(), size=9.5, color=MUTED, bold=True, first=True, font=FSB)
hline(s, ML, ty + 0.32, CW, NAVY, 1.4)
ry = ty + 0.46
for sys_name, year, approach, needs, r5, bench in rows:
    tf = tb(s, ML, ry, 2.05, 0.72)
    para(tf, sys_name, size=12, color=INK, bold=True, first=True, font=FSB)
    para(tf, year, size=8.5, color=MUTED, before=1)
    tf = tb(s, 2.85, ry, 4.75, 0.72)
    para(tf, approach, size=10.5, color=BODY, first=True, leading=1.1)
    tf = tb(s, 7.78, ry, 2.55, 0.72)
    para(tf, needs, size=10.5, color=BODY, first=True, leading=1.1)
    tf = tb(s, 10.55, ry, 2.15, 0.72)
    para(tf, [(r5, {"bold": True, "color": NAVY, "size": 13.5, "font": FSB}),
              ("   " + bench, {"size": 8.5, "color": MUTED})], first=True)
    ry += 0.80
    if sys_name != "GaussianGraph":
        hline(s, ML, ry - 0.11, CW, PANEL_LN, 0.75)
tf = tb(s, ML, ry + 0.10, CW, 0.6)
para(tf, [("Whether the geometry already present in a splat is sufficient on its own has not been "
           "systematically tested.  ",
           {"bold": True, "color": ACCENT, "size": 12, "font": FSB}),
          ("Accuracy figures as published by the respective authors, under matched evaluation "
           "protocols; not re-evaluated here.",
           {"size": 9, "color": MUTED})], first=True)

# ======================================================================
# SLIDE 6 — THE KEY IDEA
# ======================================================================
s = std_slide("The Key Idea",
              "Treat spatial relations as what they are: geometric facts.")
rect(s, ML, 1.95, CW, 1.28, PANEL, PANEL_LN, radius=0.06)
tf = tb(s, ML + 0.35, 1.95, CW - 0.7, 1.28, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("on_top_of", {"font": MONO, "color": NAVY, "bold": True}),
          (", ", {}),
          ("left_of", {"font": MONO, "color": NAVY, "bold": True}),
          (", ", {}),
          ("higher_than", {"font": MONO, "color": NAVY, "bold": True}),
          ("  —  each is fully determined by object centroids and bounding volumes, "
           "independent of what the objects are. No semantic grounding is needed to "
           "state them; none should be needed to predict them.", {"color": INK})],
     size=14, first=True, leading=1.25)
idea_cards = [
    ("LEARN", "GeoKANRelationGNN",
     "A graph network whose classification heads carry a learnable geometric metric, "
     "predicting ten relation types from splat geometry alone — 0.001 B parameters, "
     "against baselines at ~1.5 B."),
    ("SUPERVISE", "Rule-based label injection",
     "Where human annotation is missing, geometry itself supplies dense training labels, "
     "each weighted by the confidence its rule deserves."),
    ("GUARANTEE", "SceneGraphRepair",
     "Universal physical constraints repair the predicted graph deterministically — "
     "zero parameters, no training data."),
]
cx = ML
for tag, name, body in idea_cards:
    rect(s, cx, 3.55, 3.83, 2.55, WHITE, PANEL_LN, radius=0.05)
    rect(s, cx, 3.55, 3.83, 0.055, ACCENT)
    tf = tb(s, cx + 0.28, 3.80, 3.27, 2.2)
    para(tf, tag, size=10, color=ACCENT, bold=True, first=True, font=FSB, after=2)
    para(tf, name, size=14, color=INK, bold=True, font=FSB, after=6)
    para(tf, body, size=10.5, color=BODY, leading=1.2)
    cx += 3.83 + 0.30
tf = tb(s, ML, 6.35, CW, 0.4)
para(tf, "Every result in this talk tests one of these three decisions.",
     size=11.5, color=MUTED, italic=True, first=True, align=PP_ALIGN.CENTER)

# ======================================================================
# SLIDE 7 — PIPELINE
# ======================================================================
s = std_slide("System Overview",
              "From smartphone video to a logically consistent scene graph — no foundation model at inference.")
img_w = 10.9
img_h = img_w / 2.352
pic(s, FIG + r"\fig1_pipeline.png", (SW - img_w) / 2, 1.72, w=img_w)
tf = tb(s, ML, 1.72 + img_h + 0.18, CW, 0.4)
para(tf, [("Reconstruction ", {"bold": True, "color": INK}),
          ("30,000 iterations  →  ", {}),
          ("Cleaning ", {"bold": True, "color": INK}),
          ("opacity, outliers, plane  →  ", {}),
          ("HDBSCAN ", {"bold": True, "color": INK}),
          ("clustering  →  ", {}),
          ("Features ", {"bold": True, "color": INK}),
          ("10-D node · 22-D edge  →  ", {}),
          ("GeoKANRelationGNN ", {"bold": True, "color": INK}),
          ("→  ", {}),
          ("SceneGraphRepair", {"bold": True, "color": ACCENT}),
          ], size=10, color=BODY, first=True, align=PP_ALIGN.CENTER, leading=1.15)

# ======================================================================
# SLIDE 6 — GEOKAN LAYER
# ======================================================================
s = std_slide("Method  ·  Learning",
              "A GeoKAN layer learns how much each dimension matters before expanding it.")
tf = tb(s, ML, 1.90, 5.85, 2.0)
para(tf, [("A standard basis expansion treats every input dimension as equally informative. "
           "A GeoKAN layer first learns a per-dimension scaling — a ", {}),
          ("diagonal Riemannian metric g", {"bold": True, "color": INK}),
          (" — that stretches the dimensions the network has found useful and compresses "
           "the ones it has not, before any basis function is applied.", {})],
     size=12.5, color=BODY, first=True, leading=1.22)
rect(s, ML, 3.42, 5.85, 0.95, PANEL, PANEL_LN, radius=0.08)
tf = tb(s, ML + 0.3, 3.42, 5.25, 0.95, anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("z = u · √g", {"size": 20, "bold": True, "color": NAVY, "font": FSB}),
          ("      then each zᵢ is expanded against 12 fixed centres", {"size": 11, "color": MUTED})],
     first=True)
vy = 4.72
variants = [
    ("Gamma", "Input-independent metric — g = softplus(γ), one learnable scalar per dimension", "Radial basis function"),
    ("RBF", "Input-dependent metric — g = softplus(MetricNet(u)), a small two-layer network", "Radial basis function"),
    ("Wavelet", "Input-dependent metric — same MetricNet", "Mexican-hat wavelet (signed side-lobes)"),
]
tf = tb(s, ML, vy - 0.30, 5.85, 0.3)
para(tf, "THE THREE VARIANTS DIFFER IN EXACTLY TWO PLACES", size=9.5, color=MUTED,
     bold=True, first=True, font=FSB)
for name, metric, basis in variants:
    tf = tb(s, ML, vy, 5.85, 0.62)
    para(tf, [(name, {"bold": True, "color": ACCENT, "size": 12, "font": FSB}),
              ("   " + metric, {"size": 10.5, "color": BODY})], first=True, leading=1.1)
    para(tf, "Basis: " + basis, size=10, color=MUTED, before=1)
    vy += 0.68
img_w6 = 6.05
pic(s, FIG + r"\fig2_geokan_layer.png", 6.75, 1.85, w=img_w6)

# ======================================================================
# SLIDE 7 — DUAL-HEAD ARCHITECTURE
# ======================================================================
s = std_slide("Method  ·  Architecture",
              "Judging contact and judging direction call for different evidence — each receives its own head.")
lx, lw = ML, 8.55
tf = tb(s, lx, 1.88, lw, 1.35)
para(tf, [("Shared backbone.  ", {"bold": True, "color": INK}),
          ("10-D node features project to 128-D; two GATv2Conv layers (4 attention heads, residual) "
           "aggregate context, with the 22-D edge features supplied directly to the attention computation.",
           {})], size=12, color=BODY, first=True, leading=1.18, after=7)
para(tf, [("Pair representation.  ", {"bold": True, "color": INK}),
          ("p = [hᵢ, hⱼ, hᵢ−hⱼ, hᵢ⊙hⱼ] → 64-D embedding — the difference and product terms "
           "capture directional, antisymmetric structure.", {})],
     size=12, color=BODY, leading=1.18)
hy = 3.62
for title, rels, feats, dim, why, xh in [
    ("CONTACT HEAD", "on_top_of · under · attached_to · adjacent_to",
     "All 22 edge dimensions", "86-D input",
     "Needs vertical gap, containment, and contact-margin signals to judge physical support.", lx),
    ("DIRECTIONAL HEAD", "left_of · right_of · in_front_of · behind · higher_than · lower_than",
     "First 10 edge dimensions", "74-D input",
     "Direction is fully determined by relative position alone — the rest carries no signal for it.", lx + 4.4),
]:
    rect(s, xh, hy, 4.15, 2.25, PANEL, PANEL_LN, radius=0.06)
    tf = tb(s, xh + 0.25, hy + 0.2, 3.65, 1.95)
    para(tf, title, size=11, color=NAVY, bold=True, first=True, font=FSB, after=4)
    para(tf, rels, size=9.5, color=INK, font=MONO, after=5, leading=1.15)
    para(tf, [(feats, {"bold": True, "color": BODY}), ("  ·  " + dim, {"color": ACCENT, "bold": True})],
         size=10.5, after=5)
    para(tf, why, size=9.5, color=MUTED, leading=1.14)
tf = tb(s, lx, 6.14, lw, 0.6)
para(tf, [("0.001 B parameters", {"bold": True, "color": ACCENT, "size": 14, "font": FSB}),
          ("  (1,018,414 — GeoKAN-Gamma; all variants ≤ 1.1 M), against foundation-model baselines "
           "at ~1.5 B — roughly three orders of magnitude fewer.", {"size": 11, "color": BODY})],
     first=True, leading=1.15)
img_h7 = 5.15
img_w7 = img_h7 * 0.592
pic(s, FIG + r"\fig3_dualhead_architecture.png", SW - ML - img_w7 - 0.12, 1.80, h=img_h7)

# ======================================================================
# SLIDE 8 — LOGIC: INJECTION + REPAIR
# ======================================================================
s = std_slide("Method  ·  Supervision",
              "Where annotation is silent, geometry itself supplies the training signal.")
tf = tb(s, ML, 1.95, 6.9, 1.2)
para(tf, "The premise: an unannotated pair is not evidence of absence. If one object's centroid "
         "is measurably left of another's, that relation holds — whether or not a labeller "
         "recorded it.", size=12.5, color=BODY, first=True, leading=1.24)
tf = tb(s, ML, 3.15, 6.9, 3.6)
para(tf, [("Inverse completion.  ", {"bold": True, "color": INK}),
          ("Annotators record on_top_of(A,B) but rarely under(B,A); every annotation's inverse "
           "is injected at the same confidence.", {})],
     size=11.5, color=BODY, first=True, leading=1.2, after=10)
para(tf, [("Confidence-calibrated geometric rules.  ", {"bold": True, "color": INK}),
          ("Unannotated pairs receive pseudo-labels weighted by the certainty each rule "
           "deserves:", {})],
     size=11.5, color=BODY, leading=1.2, after=6)
for rel, conf, note in [
        ("Directional (left_of, higher_than, …)", "1.00", "exact, by definition"),
        ("on_top_of · under", "0.75", "high, not absolute"),
        ("attached_to", "0.60", "genuinely ambiguous"),
        ("adjacent_to", "0.55", "weakest evidence")]:
    para(tf, [(rel, {"size": 11, "color": BODY}),
              ("   " + conf, {"size": 11.5, "color": ACCENT, "bold": True, "font": FSB}),
              ("   — " + note, {"size": 9.5, "color": MUTED})], after=4)

rx = 8.15
rect(s, rx, 1.95, SW - ML - rx, 4.75, PANEL, PANEL_LN, radius=0.05)
tf = tb(s, rx + 0.32, 2.25, SW - ML - rx - 0.64, 4.2)
para(tf, "EFFECT ON THE TRAINING SET", size=10.5, color=NAVY, bold=True, first=True,
     font=FSB, after=12)
para(tf, "1,005,126", size=25, color=INK, bold=True, font=FSB, after=1)
para(tf, "directed training edges, 480 scenes", size=10, color=MUTED, after=10)
para(tf, "→  2,383,223", size=25, color=ACCENT, bold=True, font=FSB, after=1)
para(tf, "positive labels after injection", size=10, color=MUTED, after=14)
para(tf, [("Human labels are never overwritten", {"bold": True, "color": INK}),
          (" — rules apply only where annotators were silent, and each pseudo-label's "
           "confidence weights its loss.", {})],
     size=10.5, color=BODY, leading=1.2)

# ======================================================================
# SLIDE 11 — SCENEGRAPHREPAIR
# ======================================================================
s = std_slide("Method  ·  Logical Repair",
              "After inference, four universal constraints restore consistency — with zero parameters.")
tf = tb(s, ML, 1.90, CW, 0.55)
para(tf, "A neural network predicts each pair independently; physics does not work that way. "
         "SceneGraphRepair enforces what must hold in any physical scene:",
     size=12, color=BODY, first=True, leading=1.2)
cons = [
    ("Inverse completeness", "R(A,B)  ⇒  R⁻¹(B,A)",
     "the missing inverse is added, at 0.95× the source confidence"),
    ("Mutual exclusion", "on_top_of(A,B)  ∧  lower_than(A,B)  ⊥",
     "contradictory relations cannot both hold — the less confident one is removed"),
    ("Asymmetry", "R(A,B)  ∧  R(B,A)  ⊥",
     "a directional relation cannot hold both ways — the less confident one is removed"),
    ("Transitivity", "R(A,B)  ∧  R(B,C)  ⇒  R(A,C)",
     "the implied relation is added, at 0.9× the weaker link"),
]
gw, gh, ggap = 5.87, 1.62, 0.33
for i, (cname, rule, action) in enumerate(cons):
    gx = ML + (i % 2) * (gw + ggap)
    gy = 2.62 + (i // 2) * (gh + 0.26)
    rect(s, gx, gy, gw, gh, PANEL, PANEL_LN, radius=0.06)
    tf = tb(s, gx + 0.26, gy + 0.17, gw - 0.52, gh - 0.3)
    para(tf, cname.upper(), size=9.5, color=NAVY, bold=True, first=True, font=FSB, after=3)
    para(tf, rule, size=12, color=INK, bold=True, font=MONO, after=4)
    para(tf, action, size=9.5, color=MUTED, leading=1.12)
tf = tb(s, ML, 6.35, CW, 0.5)
para(tf, [("Deterministic fixed-point iteration — ", {"bold": True, "color": ACCENT, "font": FSB}),
          ("alternating removal and addition phases, at most 10 iterations; in practice every "
           "scene converges within two or three. These constraints are physical facts, not "
           "conventions of any dataset.", {"color": BODY})],
     size=11, first=True, leading=1.2)

# ======================================================================
# SLIDE 9 — DATASETS
# ======================================================================
s = std_slide("Evaluation Setup",
              "One training distribution — and three evaluation settings of increasing difficulty.")
ds_w = 9.0
ds_h = ds_w / (1356 / 518)
ds_x = (SW - ds_w) / 2
ds_y = 1.75
pic(s, FIG + r"\fig_datasets_comparison.png", ds_x, ds_y, w=ds_w)
ds_cols = [
    ("TRAINING  +  IN-DOMAIN VALIDATION",
     [("565 scenes", " — 480 train / 85 validation · 3–8 m, room scale"),
      ("26 → 10 relations", " · 1,005,126 edges → 2,383,223 labels")]),
    ("HELD OUT  —  DATASET SHIFT",
     [("4 scenes", " — ramen, teatime, waldo_kitchen, figurines"),
      ("1,644 triples", " · 1–5 m · 10–13 objects · never seen in training")]),
    ("ZERO-SHOT  —  TENFOLD SCALE SHIFT",
     [("8 scenes", " — captured for this work · 0.3–0.8 m · 4–5 objects"),
      ("508 triples", " · 3 of 8 reconstructed Z-inverted, corrected before evaluation")]),
]
col_w = ds_w / 3
for i, (role, lines) in enumerate(ds_cols):
    cx = ds_x + i * col_w + 0.12
    tf = tb(s, cx, ds_y + ds_h + 0.22, col_w - 0.24, 1.15)
    para(tf, role, size=8.5, color=ACCENT, bold=True, first=True, font=FSB, after=4)
    for val, desc in lines:
        para(tf, [(val, {"bold": True, "color": NAVY, "size": 10.5, "font": FSB}),
                  (desc, {"size": 9.5, "color": BODY})], leading=1.15, after=3)
tf = tb(s, ML, 6.60, CW, 0.35)
para(tf, "Each configuration is trained once, on 3RScan only — the held-out and tabletop "
         "evaluations reuse those weights and per-relation thresholds entirely unchanged.",
     size=10.5, color=MUTED, italic=True, first=True, align=PP_ALIGN.CENTER)

# ======================================================================
# SLIDE 10 — IN-DOMAIN RESULTS
# ======================================================================
s = std_slide("Results  ·  In-Domain",
              "In its own domain, geometry alone surpasses every foundation-model baseline.")
tf = tb(s, ML, 1.86, 6.4, 0.3)
para(tf, "PREDICATE RECALL@5, 3DSSG", size=9.5, color=MUTED, bold=True, first=True, font=FSB)
bars = [
    ("LogicSplat", "GeoKAN-Gamma · 0.001 B params · no foundation model", 98.3, ACCENT),
    ("ReLaGS", "CLIP + SAM · ~1.5 B params", 87.0, BAR_MUTE),
    ("RelationField", "~60 min per-scene optimisation, A100", 82.0, BAR_MUTE),
    ("ConceptGraphs", "large language model", 79.0, BAR_MUTE),
    ("Open3DSG", "CLIP", 65.0, BAR_MUTE),
]
by = 2.32
bar_x0 = None
for label, sub, val, fill in bars:
    bx = bar_row(s, ML, by, label, sub, val, 100.0, 3.1, fill, label_w=2.15)
    bar_x0 = bx
    by += 0.72
vline(s, bar_x0, 2.28, by - 2.28 - 0.35, PANEL_LN, 1.0)
tf = tb(s, ML, by + 0.08, 6.4, 0.5)
para(tf, "Baseline figures as published by ReLaGS — the only matched-protocol comparison "
         "covering all four systems; not re-evaluated here.", size=9, color=MUTED, first=True, leading=1.15)

rx, rw = 7.55, 5.15
tf = tb(s, rx, 1.86, rw, 0.3)
para(tf, "ARCHITECTURE COMPARISON — 85 VALIDATION SCENES", size=9.5, color=MUTED,
     bold=True, first=True, font=FSB)
tcols = [("", rx, 1.7), ("Macro F1", rx + 1.72, 1.1), ("R@5", rx + 2.86, 1.0), ("Params", rx + 3.9, 1.25)]
ty = 2.26
for name, x, w in tcols:
    tf = tb(s, x, ty, w, 0.25)
    para(tf, name, size=9.5, color=MUTED, bold=True, first=True, font=FSB)
hline(s, rx, ty + 0.28, rw, NAVY, 1.2)
trows = [
    ("MLP baseline", "0.9325", "0.9817", "1,570,666", 0),
    ("GeoKAN-Gamma", "0.9257", "0.9832", "1,018,414", 1),
    ("GeoKAN-RBF", "0.9194", "0.9818", "1,071,918", 2),
    ("GeoKAN-Wavelet", "0.9135", "0.9795", "1,071,914", 3),
]
ry = ty + 0.40
for name, f1, r5, params, i in trows:
    tf = tb(s, rx, ry, 1.72, 0.3)
    para(tf, name, size=10.5, color=INK, bold=(i == 1), first=True)
    tf = tb(s, rx + 1.72, ry, 1.1, 0.3)
    para(tf, f1, size=10.5, color=(ACCENT if f1 == "0.9325" else BODY),
         bold=(f1 == "0.9325"), first=True)
    tf = tb(s, rx + 2.86, ry, 1.0, 0.3)
    para(tf, r5, size=10.5, color=(ACCENT if r5 == "0.9832" else BODY),
         bold=(r5 == "0.9832"), first=True)
    tf = tb(s, rx + 3.9, ry, 1.25, 0.3)
    para(tf, params, size=10.5, color=BODY, first=True)
    ry += 0.44
    hline(s, rx, ry - 0.08, rw, PANEL_LN, 0.75)
tf = tb(s, rx, ry + 0.12, rw, 2.0)
para(tf, [("Macro F1 = 0.926 for the headline system.  ", {"bold": True, "color": INK})],
     size=11, first=True, after=5)
para(tf, "In-distribution the matched MLP and GeoKAN are nearly indistinguishable — the strongest "
         "single result is a property of the features, the injected labels, and the logic, not of "
         "the classification head.", size=10.5, color=BODY, leading=1.18, after=5)
para(tf, "The choice of head becomes decisive only under domain shift — the two settings "
         "that follow.", size=10.5, color=ACCENT, bold=True, leading=1.15)

# ======================================================================
# SLIDE 11 — TABLETOP RESULTS
# ======================================================================
s = std_slide("Results  ·  Scale Shift",
              "At one-tenth the training scale, zero-shot, accuracy largely holds: 92.7% Recall@5.")
lx = ML
tf = tb(s, lx, 1.86, 6.5, 0.3)
para(tf, "MICRO F1 BEFORE → AFTER SYMBOLIC REPAIR  ·  RECALL@5", size=9.5, color=MUTED,
     bold=True, first=True, font=FSB)
tcols = [("", lx, 1.85), ("Before", lx + 1.95, 0.95), ("After", lx + 2.95, 0.95),
         ("Δ", lx + 3.9, 0.8), ("R@5", lx + 4.75, 0.9)]
ty = 2.28
for name, x, w in tcols:
    tf = tb(s, x, ty, w, 0.25)
    para(tf, name, size=9.5, color=MUTED, bold=True, first=True, font=FSB)
hline(s, lx, ty + 0.28, 5.75, NAVY, 1.2)
trows = [
    ("GeoKAN-Wavelet", "0.783", "0.819", "+0.036", "93.3%", True),
    ("GeoKAN-RBF", "0.777", "0.802", "+0.025", "94.3%", False),
    ("GeoKAN-Gamma", "0.740", "0.774", "+0.034", "92.7%", False),
    ("MLP baseline", "0.714", "0.721", "+0.007", "88.6%", False),
]
ry = ty + 0.42
for name, b, a, d, r5, best in trows:
    tf = tb(s, lx, ry, 1.95, 0.3)
    para(tf, name, size=11, color=INK, bold=best, first=True)
    tf = tb(s, lx + 1.95, ry, 0.95, 0.3)
    para(tf, b, size=11, color=BODY, first=True)
    tf = tb(s, lx + 2.95, ry, 0.95, 0.3)
    para(tf, a, size=11, color=INK, bold=True, first=True)
    tf = tb(s, lx + 3.9, ry, 0.8, 0.3)
    para(tf, d, size=11, color=ACCENT, bold=True, first=True)
    tf = tb(s, lx + 4.75, ry, 0.9, 0.3)
    para(tf, r5, size=11, color=NAVY, bold=True, first=True)
    ry += 0.5
    hline(s, lx, ry - 0.11, 5.75, PANEL_LN, 0.75)
tf = tb(s, lx, ry + 0.10, 5.9, 1.6)
para(tf, "All four are LogicSplat configurations trained on 3RScan only — no tabletop fine-tuning, "
         "thresholds carried over unmodified. 8 scenes, 508 ground-truth triples. Recall@3/5 are "
         "computed from raw rankings and are unchanged by repair.",
     size=9.5, color=MUTED, first=True, leading=1.18)

rx = 7.05
rect(s, rx, 1.90, SW - ML - rx, 4.85, PANEL, PANEL_LN, radius=0.05)
tf = tb(s, rx + 0.3, 2.14, SW - ML - rx - 0.6, 4.4)
para(tf, "WHAT THE SCALE SHIFT REVEALS", size=10.5, color=NAVY, bold=True, first=True,
     font=FSB, after=8)
para(tf, [("Every GeoKAN variant beats the matched MLP on every column ",
           {"bold": True, "color": INK}),
          ("— by 5.3 to 9.8 Micro F1 points and 4.1 to 5.7 Recall@5 points after repair.", {})],
     size=11.5, color=BODY, leading=1.18, after=9)
para(tf, [("The in-distribution ranking inverts. ", {"bold": True, "color": INK}),
          ("Wavelet — weakest in-domain — leads cross-domain: the input-conditioned metric "
           "generalises further than Gamma's fixed scalar.", {})],
     size=11.5, color=BODY, leading=1.18, after=9)
para(tf, [("Repair adds, almost never removes. ", {"bold": True, "color": INK}),
          ("88 / 54 / 42 relations added (Gamma / RBF / Wavelet), only 0 / 8 / 0 removed — "
           "every variant converges within 2 iterations.", {})],
     size=11.5, color=BODY, leading=1.18)

# ======================================================================
# SLIDE 12 — LERF RESULTS
# ======================================================================
s = std_slide("Results  ·  Dataset Shift",
              "On four held-out LERF scenes: 95.0% Recall@5, 31.8 points above GaussianGraph's strongest configuration.")
tf = tb(s, ML, 1.86, 7.0, 0.3)
para(tf, "RECALL@5 — POSITIONAL QUERIES, 1,644 GROUND-TRUTH TRIPLES", size=9.5,
     color=MUTED, bold=True, first=True, font=FSB)
bars = [
    ("LogicSplat · GeoKAN-Gamma", "after repair", 95.0, ACCENT),
    ("LogicSplat · MLP baseline", "after repair", 92.8, BAR_NAVY),
    ("GaussianGraph", "LLaVA + 3D correction — its best configuration", 63.2, BAR_MUTE),
    ("GaussianGraph", "LLaVA, no 3D correction", 49.8, BAR_MUTE),
]
by = 2.32
for label, sub, val, fill in bars:
    bx = bar_row(s, ML, by, label, sub, val, 100.0, 3.15, fill, label_w=2.75)
    by += 0.78
vline(s, bx, 2.28, by - 2.28 - 0.40, PANEL_LN, 1.0)
tf = tb(s, ML, by + 0.05, 7.0, 0.5)
para(tf, "GaussianGraph figures as published for positional queries — the relation types "
         "evaluated here are entirely positional.", size=9, color=MUTED,
     first=True, leading=1.15)

rx = 8.0
rect(s, rx, 1.90, SW - ML - rx, 4.85, PANEL, PANEL_LN, radius=0.05)
tf = tb(s, rx + 0.3, 2.14, SW - ML - rx - 0.6, 4.4)
para(tf, "REPAIR ON A DATASET SHIFT", size=10.5, color=NAVY, bold=True, first=True, font=FSB, after=8)
para(tf, [("Recall@3:  71.5% → 89.9%", {"bold": True, "color": ACCENT, "size": 15, "font": FSB})],
     first=False, after=2)
para(tf, "+18.4 points from symbolic repair alone — no retraining.", size=10.5, color=BODY,
     leading=1.15, after=10)
para(tf, [("Micro F1 = 90.0%", {"bold": True, "color": INK, "size": 13, "font": FSB}),
          ("   after repair", {"size": 10, "color": MUTED})], after=2)
para(tf, "Repair adds 39 relations and removes none — the predictions stay internally "
         "consistent even off-distribution.", size=10.5, color=BODY, leading=1.16, after=10)
para(tf, [("Per scene (R@5 after):  ", {"bold": True, "color": INK, "size": 10.5})], after=2)
para(tf, "ramen 99.0% · waldo_kitchen 93.7% · figurines 92.3% · teatime 91.5%",
     size=10.5, color=BODY, leading=1.3)

# ======================================================================
# SLIDE 13 — QUALITATIVE
# ======================================================================
s = std_slide("Results  ·  Qualitative",
              "A complete pass over one real scene: 66 of 84 predicted relations verified correct.")
imh = 3.55
w1 = imh * 1.638   # 5.82
w2 = imh * 1.4184  # true aspect of the splat figure
w3 = imh * 0.98    # 3.48
total = w1 + w2 + w3
scale = min(1.0, (CW - 0.5) / total)
imh *= scale
w1 *= scale; w2 *= scale; w3 *= scale
gap = (CW - (w1 + w2 + w3)) / 2
x1, x2, x3 = ML, ML + w1 + gap, ML + w1 + gap + w2 + gap
iy = 2.25
for x, w, step in [(x1, w1, "1  ·  SMARTPHONE CAPTURE"),
                   (x2, w2, "2  ·  CLUSTERED SPLAT (COLOUR = CLUSTER)"),
                   (x3, w3, "3  ·  PREDICTED GRAPH, POST-REPAIR")]:
    tf = tb(s, x, iy - 0.32, w, 0.3)
    para(tf, step, size=9, color=ACCENT, bold=True, first=True, font=FSB)
pic(s, FIG + r"\fig_input_photo.png", x1, iy, w=w1)
pic(s, FIG + r"\fig_splat_reconstruction.png", x2, iy, w=w2)
pic(s, FIG + r"\fig5_qualitative_example.png", x3, iy, w=w3)
tf = tb(s, ML, iy + imh + 0.18, CW, 0.7)
para(tf, [("Router, hair-dryer box, water bottle, watch, and pen.  ", {"bold": True, "color": INK}),
          ("The router, box, watch, and bottle sit so close in raw position that colour is what "
           "separates them at the clustering step — the graph shown is the repaired prediction, "
           "restricted to a verified-correct subset of relations.", {})],
     size=11, color=BODY, first=True, leading=1.2, align=PP_ALIGN.CENTER)

# ======================================================================
# SLIDE 13B — QUALITATIVE, LERF (dataset-shift counterpart)
# ======================================================================
s = std_slide("Results  ·  Qualitative  ·  LERF",
              "The same pass on a public benchmark, zero-shot: 198 of 228 predicted relations verified correct.")
imh = 3.55
w1 = imh * (988 / 730)    # real aspect of lerf_teatime.jpg
w2 = imh * 1.0            # square top-down splat panel
w3 = imh * 0.98
total = w1 + w2 + w3
scale = min(1.0, (CW - 0.5) / total)
imh *= scale
w1 *= scale; w2 *= scale; w3 *= scale
gap = (CW - (w1 + w2 + w3)) / 2
x1, x2, x3 = ML, ML + w1 + gap, ML + w1 + gap + w2 + gap
iy = 2.25
for x, w, step in [(x1, w1, "1  ·  PUBLIC BENCHMARK PHOTO (TEATIME)"),
                   (x2, w2, "2  ·  CLUSTERED SPLAT (COLOUR = CLUSTER)"),
                   (x3, w3, "3  ·  PREDICTED GRAPH, POST-REPAIR")]:
    tf = tb(s, x, iy - 0.32, w, 0.3)
    para(tf, step, size=9, color=ACCENT, bold=True, first=True, font=FSB)
pic(s, r"demo\data\lerf_teatime.jpg", x1, iy, w=w1)
pic(s, FIG + r"\fig_splat_reconstruction_lerf_teatime.png", x2, iy, w=w2)
pic(s, FIG + r"\fig_qualitative_example_lerf_teatime.png", x3, iy, w=w3)
tf = tb(s, ML, iy + imh + 0.18, CW, 0.7)
para(tf, [("A stuffed bear's café table: cookies, coffee, tea, and a plush sheep.  ", {"bold": True, "color": INK}),
          ("The splat is seen top-down, oriented to match the photo — the plate and its three cookies "
           "are visible at centre. Object names are the benchmark's own annotations (the sheep is "
           "labelled “hooves”). Same checkpoint, no fine-tuning, thresholds unchanged.", {})],
     size=11, color=BODY, first=True, leading=1.2, align=PP_ALIGN.CENTER)

# ======================================================================
# SLIDE 14 — GEOKAN vs MLP ABLATION
# ======================================================================
s = std_slide("Analysis  ·  Ablation",
              "GeoKAN matches the MLP in-distribution — and consistently surpasses it under domain shift.")
img_w14 = 5.55
pic(s, FIG + r"\fig4_mlp_ablation.png", ML, 1.95, w=img_w14)
rx = ML + img_w14 + 0.45
rw = SW - ML - rx
tf = tb(s, rx, 1.95, rw, 4.9)
para(tf, [("Controlled ablation. ", {"bold": True, "color": INK}),
          ("Only the layer changes — same features, labels, augmentation, loss, schedule, "
           "thresholds. The MLP gets more capacity, not less: 1,570,666 vs 1,018,414 parameters.", {})],
     size=11.5, color=BODY, first=True, leading=1.2, after=10)
para(tf, [("In-distribution — no advantage.  ", {"bold": True, "color": INK}),
          ("MLP Macro F1 0.9325 vs GeoKAN-Gamma 0.9257; consistent with Yu et al.'s "
           "KAN-vs-MLP findings.", {})], size=11.5, color=BODY, leading=1.2, after=8)
para(tf, [("Scale shift (tabletop) — GeoKAN leads.  ", {"bold": True, "color": ACCENT}),
          ("Every variant beats the MLP on every metric: +5.3 to +9.8 Micro F1, "
           "+4.1 to +5.7 Recall@5 after repair.", {})], size=11.5, color=BODY, leading=1.2, after=8)
para(tf, [("Dataset shift (LERF) — GeoKAN leads.  ", {"bold": True, "color": ACCENT}),
          ("+2.2 points in both Micro F1 and Recall@5 — same direction, independent setting.", {})],
     size=11.5, color=BODY, leading=1.2, after=12)
para(tf, "The metric warp is an inductive bias that costs nothing in-distribution — and repays "
         "that cost precisely where the system is meant to operate: beyond its training domain.",
     size=12, color=NAVY, italic=True, leading=1.2)

# ======================================================================
# SLIDE 15 — CONTRIBUTIONS
# ======================================================================
s = std_slide("Contributions",
              "Five findings, each established by a controlled comparison.")
contribs = [
    ("C1", "Semantic grounding is not a prerequisite",
     "Geometry-only prediction from Gaussian-splat geometry exceeds every compared foundation-model "
     "or per-scene-optimised method on 3DSSG."),
    ("C2", "Annotation sparsity has a principled fix",
     "Rule-based label injection converts geometrically certain unannotated pairs into pseudo-labels "
     "calibrated to each rule's actual certainty — instead of false negatives."),
    ("C3", "GeoKAN buys generalisation, not accuracy",
     "No in-distribution advantage over a matched MLP, but a consistent gain under both scale shift "
     "and dataset shift — a domain-shift axis the KAN-vs-MLP literature had not tested."),
    ("C4", "Symbolic repair works by recall completion",
     "Across two held-out settings its contribution is almost entirely completing under-predicted "
     "relations — out-of-distribution errors are dominated by omission, not self-contradiction."),
    ("C5", "Scale, not dataset identity, drives degradation",
     "Comparing the two shifts isolates scale as the dominant factor in cross-domain transfer; "
     "in_front_of and behind degrade most, on_top_of and lower_than next."),
]
oy = 1.92
for tag, title, body in contribs:
    tf = tb(s, ML, oy, 0.75, 0.5)
    para(tf, tag, size=17, color=ACCENT, bold=True, first=True, font=FSB)
    tf = tb(s, ML + 0.85, oy, CW - 0.85, 0.95)
    para(tf, [(title + ".   ", {"bold": True, "color": INK, "size": 12.5, "font": FSB}),
              (body, {"size": 11, "color": BODY})], first=True, leading=1.16)
    oy += 0.99
    if tag != "C5":
        hline(s, ML + 0.85, oy - 0.14, CW - 0.85, PANEL_LN, 0.75)

# ======================================================================
# SLIDE 16 — LIMITATIONS & FUTURE WORK
# ======================================================================
s = std_slide("Limitations & Future Work",
              "Every limitation here is a deliberate scope choice — and each defines a concrete next step.")
lx, lw = ML, 5.87
tf = tb(s, lx, 1.92, lw, 0.3)
para(tf, "SCOPE CHOICES IN THIS STUDY", size=10.5, color=NAVY, bold=True, first=True, font=FSB)
tf = tb(s, lx, 2.30, lw, 4.5)
lims = [
    ("Oracle clustering for held-out settings", "the standard PredCls convention — isolates relation "
     "prediction, this study's contribution, from class-agnostic instance segmentation."),
    ("Scenes span 4–13 objects", "comfortably covering the tabletop-manipulation and single-room "
     "scenarios this system targets."),
    ("attached_to has no cross-domain test", "it never occurs in the tabletop benchmark's natural "
     "arrangements."),
    ("Single training run per configuration", "consistent with a first systematic comparison across "
     "architectures and settings."),
    ("Static, rigid, indoor scenes", "the geometric assumptions are designed for the deployment "
     "scenarios motivating this work."),
]
first = True
for head, body in lims:
    para(tf, [(head + " — ", {"bold": True, "color": INK}), (body, {})],
         size=10.5, color=BODY, leading=1.16, after=8, first=first)
    first = False

rx = ML + lw + 0.33
rect(s, rx, 1.90, lw, 4.85, PANEL, PANEL_LN, radius=0.05)
tf = tb(s, rx + 0.3, 2.14, lw - 0.6, 4.4)
para(tf, "WHERE THIS GOES NEXT", size=10.5, color=NAVY, bold=True, first=True, font=FSB, after=9)
futs = [
    ("Exploit the splat's own covariance", "planarity, elongation, sphericity are already encoded "
     "per-Gaussian but unused by the current feature set."),
    ("Denser, multi-room environments", "extend evaluation beyond 13 objects toward broader "
     "real-world deployment."),
    ("attached_to cross-domain transfer", "with scenes where it naturally occurs."),
    ("Multiple seeds for the closest margins", "particularly the in-distribution gap between "
     "GeoKAN-Gamma and the matched MLP."),
    ("The generalisation question, beyond this task", "test whether GeoKAN's advantage under "
     "distribution shift extends to other KAN applications."),
]
for head, body in futs:
    para(tf, [(head + " — ", {"bold": True, "color": INK}), (body, {})],
         size=10.5, color=BODY, leading=1.16, after=8)

# ======================================================================
# SLIDE 21B — RECORDED PIPELINE RUN (embedded video)
# ======================================================================
s = std_slide("Demonstration  ·  Recorded Run",
              "The complete pipeline, captured in one unedited recording.")
VID = r"C:\Users\Debjyoti\Desktop\LogicSplat\media\pipeline_recording.mp4"
VID_POSTER = r"C:\Users\Debjyoti\Desktop\LogicSplat\media\pipeline_recording_poster.png"
vid_w = 8.2
vid_h = vid_w / (1908 / 956)   # true aspect of the recording
vx = (SW - vid_w) / 2
vy = 1.80
rect(s, vx - 0.03, vy - 0.03, vid_w + 0.06, vid_h + 0.06, None, PANEL_LN, 1.0)
s.shapes.add_movie(VID, Inches(vx), Inches(vy), Inches(vid_w), Inches(vid_h),
                   poster_frame_image=VID_POSTER, mime_type="video/mp4")
tf = tb(s, ML, vy + vid_h + 0.12, CW, 0.3)
para(tf, "Click to play — a 19-minute unedited screen recording, embedded in this file.",
     size=10, color=MUTED, italic=True, first=True, align=PP_ALIGN.CENTER)
vch_w = (CW - 2 * 0.3) / 3
vch_y = vy + vid_h + 0.52
for i, (num, txt) in enumerate([
        ("1", "RECONSTRUCTION — COLMAP + splatfacto, 30,000 iterations"),
        ("2", "PIPELINE — clustering → GeoKANRelationGNN → SceneGraphRepair"),
        ("3", "LOGICSPLAT LIVE — every prediction inspected interactively")]):
    cx = ML + i * (vch_w + 0.3)
    rect(s, cx, vch_y, vch_w, 0.5, PANEL, PANEL_LN, radius=0.14)
    tf = tb(s, cx + 0.2, vch_y, vch_w - 0.4, 0.5, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(num + "  ·  ", {"bold": True, "color": ACCENT, "font": FSB}),
              (txt, {"color": INK})], size=9, first=True, bold=True, font=FSB)

# ======================================================================
# SLIDE 17 — CONCLUSION / THANK YOU
# ======================================================================
s = slide_new(bg=NAVY)
pts2 = [(11.0, 0.9), (12.3, 1.6), (10.4, 2.1), (12.0, 3.0)]
motif(s, pts2, [(0, 1), (0, 2), (1, 3), (2, 3)], accent_idx=1, r=0.09)
tf = tb(s, ML, 1.15, 11.0, 0.9)
para(tf, "Geometry, a learnable metric, and deterministic logic suffice —",
     size=22, color=CLOUD, first=True, leading=1.1)
para(tf, [("no foundation model required.", {"color": WHITE, "bold": True, "size": 22, "font": FSB})],
     leading=1.1)
sy = 2.75
stats = [
    ("98.3%", "Recall@5\nin-domain (3DSSG)"),
    ("95.0%", "Recall@5\nheld-out (LERF)"),
    ("92.7%", "Recall@5\nzero-shot tabletop"),
    ("0.001 B", "parameters (1.02 M)\nbaselines: ~1.5 B"),
    ("5–6 s", "end-to-end\nper scene"),
]
chip_w = 2.25
gap = (CW - 5 * chip_w) / 4
cx = ML
for val, lab in stats:
    rect(s, cx, sy, chip_w, 1.55, NAVY_2, radius=0.08)
    tf = tb(s, cx, sy + 0.22, chip_w, 0.6)
    para(tf, val, size=26, color=ACCENT_D, bold=True, align=PP_ALIGN.CENTER, first=True, font=FSB)
    tf = tb(s, cx, sy + 0.88, chip_w, 0.6)
    for i, line in enumerate(lab.split("\n")):
        para(tf, line, size=9.5, color=CLOUD, align=PP_ALIGN.CENTER, first=(i == 0), leading=1.1)
    cx += chip_w + gap
tf = tb(s, ML, 5.1, CW, 0.9)
para(tf, "Thank you.", size=40, color=WHITE, bold=True, first=True, font=FSB)
tf = tb(s, ML, 6.05, CW, 0.4)
para(tf, "Questions & discussion welcome.", size=15, color=CLOUD, first=True)
tf = tb(s, ML, 6.75, CW, 0.4)
para(tf, "Debjyoti Sengupta  ·  debjyotiashu66@gmail.com  ·  Symbiosis Institute of Geoinformatics",
     size=10.5, color=FAINT, first=True)

prs.save(OUT)
print("Saved", OUT, "slides:", len(prs.slides.__iter__.__self__._sldIdLst))
